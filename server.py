from flask import Flask, request, jsonify
from pptx import Presentation
import warnings
from urllib3.exceptions import InsecureRequestWarning

# Initialize the Flask application
app = Flask(__name__)

# Suppress the specific warning
warnings.simplefilter('ignore', InsecureRequestWarning)

@app.route('/api/data', methods=['POST'])
def receive_data_and_create_ppt():
    # Extract JSON data from the incoming request
    data = request.json
    
    # Check if data contains the required fields
    if data and "company_name" in data and "revenue" in data and "tech_stack" in data:
        # Path to your template PowerPoint file
        pptx_file = '/Users/deepaksarkar/Documents/pptascode/catalyst.pptx'
        prs = Presentation(pptx_file)

        # Iterate through placeholders and replace them with data
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "{{company_name}}" in shape.text:
                        shape.text = shape.text.replace("{{company_name}}", data["company_name"])
                    if "{{revenue}}" in shape.text:
                        shape.text = shape.text.replace("{{revenue}}", data["revenue"])
                    if "{{tech_stack}}" in shape.text:
                        tech_stack_str = ', '.join(data["tech_stack"])
                        shape.text = shape.text.replace("{{tech_stack}}", tech_stack_str)

        # Save the populated presentation
        output_file = '/Users/deepaksarkar/Documents/pptascode/populated_presentation.pptx'
        prs.save(output_file)

        # Return a success response
        return jsonify({"message": "Presentation successfully created", "output_file": output_file}), 200
    else:
        # Return an error response if required data is missing
        return jsonify({"error": "Missing data in request"}), 400

if __name__ == '__main__':
    # Run the Flask app on port 5000
    app.run(debug=True, port=5000)
